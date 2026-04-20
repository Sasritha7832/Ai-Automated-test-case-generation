from pptx import Presentation

def add_slide(prs, title, content_bullets):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title
    
    tf = body_shape.text_frame
    if content_bullets:
        tf.text = content_bullets[0]
        for bullet in content_bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
        
def create_ppt():
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Intelligent QA Automation Platform"
    subtitle.text = "AI-Powered Test Case Generation from PRD to Production\n\nPresented by: [Member 1], [Member 2], [Member 3], [Member 4]"

    # Slide 2: The QA Problem
    add_slide(prs, "The Software Testing Bottleneck", [
        "Time Consuming: Reading a 50+ page PRD takes days.",
        "Manual Effort: Writing hundreds of test cases by hand is slow.",
        "Human Error: Testers frequently miss edge cases and complex security flaws.",
        "Redundancy: High chance of writing duplicate test scenarios."
    ])

    # Slide 3: Our Solution
    add_slide(prs, "Automating QA with Artificial Intelligence", [
        "Instant Analysis: Converts massive PRD PDFs into test suites in seconds.",
        "Deep Coverage: Generates Functional, Security, and UI test cases.",
        "Smart Priority: Uses Machine Learning to rank bugs by severity.",
        "Seamless Export: One-click download to Jira, TestRail, and Excel."
    ])

    # Slide 4: Platform Architecture
    add_slide(prs, "How It Works (The 4 Pillars)", [
        "1. Document Parsing (RAG): Extracts and embeds text logic.",
        "2. Generative AI: Reasons like a Senior QA Lead to write tests.",
        "3. Data Science (ML): Predicts risk and scores test priorities.",
        "4. Computer Vision: Scans raw UI screenshots for visual testing."
    ])

    # Slide 5: Module 1
    add_slide(prs, "Module 1 - Teaching the AI the Rules (RAG)", [
        "The Tech: LangChain + Retrieval-Augmented Generation (RAG).",
        "Chunking: Slices the PRD into logical, mathematical text blocks.",
        "Embedding Base: Uses all-MiniLM-L6-v2 to convert text into vectors.",
        "FAISS Database: Instantly searches and retrieves only the exact rules needed for a test."
    ])

    # Slide 6: Module 2
    add_slide(prs, "Module 2 - The Core AI Engine", [
        "The Intelligence: Powered by LLaMA 3 (8B) via Groq API.",
        "Strict Grounding: AI is physically restricted to the FAISS database to prevent hallucination.",
        "Deduplication Engine: Uses TF-IDF and Cosine Similarity to automatically merge duplicate test cases.",
        "Output: Highly structured JSON consisting of Steps, Expected Results, and Preconditions."
    ])

    # Slide 7: Module 3
    add_slide(prs, "Module 3 - Predicting Risks & Priorities", [
        "The Model: Random Forest Machine Learning Classifier.",
        "Feature Engineering: Scans test complexity, step count, and high-risk NLP keywords.",
        "Scoring: Automatically flags tests from Priority 0 (Critical) to Priority 3 (Low).",
        "Result: Forces developers to focus on fixing the most dangerous bugs first."
    ])

    # Slide 8: Module 4
    add_slide(prs, "Module 4 - Testing What Users Actually See", [
        "The Tech: OpenCV + Tesseract OCR Engine.",
        "Image Processing: Applies Gaussian Blur and Canny Edge Detection to user-uploaded screenshots.",
        "Detection: Mathematically isolates bounding boxes around UI buttons and text fields.",
        "Result: Generates visual test cases (e.g., 'Verify the blue Submit button is clickable')."
    ])

    # Slide 9: QA Chatbot
    add_slide(prs, "Interactive QA Chatbot", [
        "Built-in Assistant: An AI chatbot deeply integrated into the platform.",
        "Context-Aware: Reads both the uploaded PRD and the generated test cases simultaneously.",
        "Use Case: Developers can ask complex questions like 'What happens if a user enters the wrong password?'",
        "Accuracy: Returns exact, technical answers grounded purely in the document's rules."
    ])

    # Slide 10: Analytics Dashboard
    add_slide(prs, "Real-Time Coverage Tracking", [
        "Visualizing Data: Built natively using Plotly Express.",
        "Coverage Heat Maps: Shows exact PRD Requirement IDs vs. Test Scenario coverage density.",
        "Scenario Graphs: Maps the workflow of critical test cases.",
        "Impact: Allows a QA Manager to instantly spot missing code coverage."
    ])

    # Slide 11: Business Impact
    add_slide(prs, "Why This Platform Matters", [
        "Speed: 5 days of manual reading reduced to 30 seconds of compute.",
        "Cost Efficiency: Saves companies thousands of dollars in wasted QA engineering hours.",
        "Quality Assurance: The AI never gets tired and catches maximum edge-case scenarios."
    ])

    # Slide 12: Conclusion
    add_slide(prs, "Thank You", [
        "Code Repository: https://github.com/Sasritha7832/Ai-Automated-test-case-generation",
        "Are there any questions?",
        "[Optional: Live Demonstration Uploading an Enterprise PRD]"
    ])

    prs.save('Intelligent_QA_Platform_Presentation.pptx')
    print("Presentation generated successfully!")

if __name__ == '__main__':
    create_ppt()
